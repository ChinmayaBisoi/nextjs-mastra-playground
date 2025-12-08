from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import json
from typing import Dict, Any, List

app = FastAPI()

# CORS middleware - configure origins in production
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # TODO: Restrict to your Vercel domain in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def extract_shape_data(shape) -> Dict[str, Any]:
    """Extract data from a shape object."""
    shape_data: Dict[str, Any] = {
        "shape_id": shape.shape_id,
        "shape_type": str(shape.shape_type),
        "name": shape.name if hasattr(shape, "name") else None,
    }

    # Extract text if available
    if hasattr(shape, "text"):
        shape_data["text"] = shape.text

    # Extract text frame details
    if hasattr(shape, "text_frame"):
        text_frame = shape.text_frame
        shape_data["text_frame"] = {
            "has_text": text_frame.has_text,
            "paragraphs": [],
        }
        for paragraph in text_frame.paragraphs:
            para_data = {
                "text": paragraph.text,
                "level": paragraph.level,
                "alignment": str(paragraph.alignment) if paragraph.alignment else None,
            }
            # Extract runs (formatting)
            para_data["runs"] = []
            for run in paragraph.runs:
                run_data = {
                    "text": run.text,
                    "bold": run.font.bold if run.font.bold is not None else False,
                    "italic": run.font.italic if run.font.italic is not None else False,
                    "underline": run.font.underline if run.font.underline is not None else False,
                }
                if run.font.size:
                    run_data["font_size"] = str(run.font.size)
                if run.font.name:
                    run_data["font_name"] = run.font.name
                if run.font.color:
                    run_data["font_color"] = str(run.font.color.type)
                para_data["runs"].append(run_data)
            shape_data["text_frame"]["paragraphs"].append(para_data)

    # Extract image if available
    if hasattr(shape, "image"):
        shape_data["has_image"] = True
        shape_data["image_format"] = shape.image.ext
        shape_data["image_size"] = {
            "width": shape.image.width,
            "height": shape.image.height,
        }

    # Extract table if available
    if hasattr(shape, "table"):
        table = shape.table
        shape_data["table"] = {
            "rows": len(table.rows),
            "columns": len(table.columns),
            "cells": [],
        }
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_data = {
                    "row": row_idx,
                    "column": col_idx,
                    "text": cell.text,
                }
                shape_data["table"]["cells"].append(cell_data)

    # Extract position and size
    if hasattr(shape, "left") and hasattr(shape, "top"):
        shape_data["position"] = {
            "left": shape.left,
            "top": shape.top,
        }
    if hasattr(shape, "width") and hasattr(shape, "height"):
        shape_data["size"] = {
            "width": shape.width,
            "height": shape.height,
        }

    return shape_data


@app.get("/")
async def root():
    return {"status": "ok", "service": "pptx-parser"}


@app.get("/health")
async def health():
    return {"status": "healthy"}


@app.post("/parse-pptx")
async def parse_pptx(file: UploadFile = File(...)):
    """
    Parse a PPTX file using python-pptx and return structured data.
    """
    if not file.filename or not file.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="File must be a .pptx file")

    try:
        # Read file content
        contents = await file.read()

        # Check file size (16MB limit)
        MAX_FILE_SIZE = 16 * 1024 * 1024  # 16MB
        if len(contents) > MAX_FILE_SIZE:
            raise HTTPException(
                status_code=400,
                detail=f"File size exceeds 16MB limit. Your file is {len(contents) / (1024 * 1024):.2f}MB",
            )

        # Parse with python-pptx
        prs = Presentation(io.BytesIO(contents))

        # Extract presentation-level metadata
        presentation_data = {
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "slide_count": len(prs.slides),
        }

        # Extract slides data
        slides_data = []
        for i, slide in enumerate(prs.slides, 1):
            slide_info: Dict[str, Any] = {
                "slide_number": i,
                "slide_id": slide.slide_id,
                "shapes": [],
            }

            # Extract shapes
            for shape in slide.shapes:
                try:
                    shape_data = extract_shape_data(shape)
                    slide_info["shapes"].append(shape_data)
                except Exception as e:
                    # Log error but continue processing other shapes
                    print(f"Error extracting shape {shape.shape_id}: {str(e)}")
                    slide_info["shapes"].append(
                        {
                            "shape_id": shape.shape_id,
                            "error": str(e),
                        }
                    )

            slides_data.append(slide_info)

        return {
            "success": True,
            "fileName": file.filename,
            "fileSize": len(contents),
            "presentation": presentation_data,
            "slides": slides_data,
        }

    except HTTPException:
        raise
    except Exception as e:
        print(f"Error parsing PPTX: {str(e)}")
        raise HTTPException(
            status_code=500, detail=f"Failed to parse PPTX file: {str(e)}"
        )

